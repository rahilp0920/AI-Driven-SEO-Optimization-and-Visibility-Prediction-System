"""Generate the CIS 2450 final-project slide deck.

A single-file builder. Reads the corpus stats and saved metrics live so the
slides always reflect the committed state. Renders into
``presentation/slides.pptx``.

Run:
    python -m presentation.build_charts        # regenerate PNGs first
    python -m presentation.build_slides
"""

from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ── Paths ──────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parents[1]
CHART_DIR = ROOT / "presentation" / "charts"
FEATURES_CSV = ROOT / "data" / "processed" / "features.csv"
METRICS_DIR = ROOT / "models" / "metrics"
OUT_PPTX = ROOT / "presentation" / "slides.pptx"

# ── Theme ──────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INDIGO = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DARK = RGBColor(0x37, 0x30, 0xA3)
ACCENT = RGBColor(0xEC, 0x48, 0x99)
SLATE = RGBColor(0x0F, 0x17, 0x2A)
SLATE_DIM = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
SOFT = RGBColor(0xEE, 0xF2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xFA, 0xFC)
GOOD = RGBColor(0x10, 0xB9, 0x81)
BAD = RGBColor(0xEF, 0x44, 0x44)
WARN = RGBColor(0xF5, 0x9E, 0x0B)


@dataclass
class Theme:
    """Centralised typography + spacing knobs."""

    font_family: str = "Inter"
    font_fallback: str = "Calibri"


THEME = Theme()


# ── Low-level helpers ──────────────────────────────────────────────────


def _set_solid_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_text(tf, text: str, *, size: int = 16, bold: bool = False,
              color: RGBColor = SLATE, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, family: str | None = None) -> None:
    """Replace text-frame contents with a single styled run."""
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = family or THEME.font_family
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_paragraphs(tf, lines: list[tuple[str, dict]],
                    base_color: RGBColor = SLATE,
                    base_size: int = 16,
                    bullet: bool = False) -> None:
    """Append multiple styled paragraphs. Each ``(text, opts)`` tuple may
    override ``size``, ``bold``, ``color``, ``align``, ``space_after``,
    ``family`` per paragraph."""
    tf.word_wrap = True
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 and not tf.paragraphs[0].text else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])
        run = p.add_run()
        run.text = ("• " + text) if bullet else text
        run.font.name = opts.get("family", THEME.font_family)
        run.font.size = Pt(opts.get("size", base_size))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = opts.get("color", base_color)


def _add_rect(slide, x, y, w, h, color: RGBColor) -> None:
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_solid_fill(s, color)
    s.shadow.inherit = False
    return s


def _add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(x, y, w, h).text_frame


# ── Shared chrome: header strip, footer, accent corner ─────────────────


def _slide_chrome(slide, *, eyebrow: str, title: str, page_idx: int,
                  total_pages: int) -> None:
    """Apply consistent chrome to every content slide."""
    # Subtle top accent strip — 6 px tall.
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), INDIGO)
    # Background block: full slide white-on-bg gradient effect via two layers.
    # python-pptx can't render true gradients, so we layer a white panel inset
    # by 0 — the slide background itself is white via the master.

    # Eyebrow tag (small, uppercase).
    eb_tf = _add_textbox(slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.4))
    _set_text(eb_tf, eyebrow.upper(), size=11, bold=True, color=INDIGO,
              family=THEME.font_family)
    p = eb_tf.paragraphs[0]
    for run in p.runs:
        run.font.name = "Inter"

    # Title — large, dark, tight.
    title_tf = _add_textbox(slide, Inches(0.55), Inches(0.78),
                            Inches(12.0), Inches(1.0))
    _set_text(title_tf, title, size=34, bold=True, color=SLATE)

    # Page number — bottom right.
    pn_tf = _add_textbox(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3))
    _set_text(pn_tf, f"{page_idx} / {total_pages}", size=10,
              color=MUTED, align=PP_ALIGN.RIGHT)

    # Footer — bottom left, project tag.
    foot_tf = _add_textbox(slide, Inches(0.6), Inches(7.05), Inches(8.0), Inches(0.3))
    _set_text(foot_tf, "SEO Ranking Predictor · CIS 2450 Final",
              size=10, color=MUTED)


def _new_slide(prs: Presentation):
    """Append a blank slide with white background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return slide


def _add_pill(slide, x, y, label: str, color: RGBColor = INDIGO,
              text_color: RGBColor = WHITE, size: int = 11) -> None:
    """Tag-style pill — used for category labels on diagrams."""
    h = Inches(0.32)
    w = Inches(max(1.0, len(label) * 0.09 + 0.4))
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    _set_solid_fill(pill, color)
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.name = THEME.font_family
    run.font.size = Pt(size); run.font.bold = True
    run.font.color.rgb = text_color


def _add_stat_card(slide, x, y, w, h, label: str, value: str,
                   value_color: RGBColor = INDIGO_DARK) -> None:
    """White card with eyebrow label and large value — used on title and
    summary slides for headline numbers."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    _set_solid_fill(card, WHITE)
    card.line.color.rgb = LINE
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.18)
    _add_paragraphs(tf, [
        (label.upper(), {"size": 10, "bold": True, "color": MUTED,
                         "space_after": 4}),
        (value, {"size": 28, "bold": True, "color": value_color}),
    ])


def _add_callout(slide, x, y, w, h, body: str) -> None:
    """Soft indigo card with a small left accent — for short takeaways."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.adjustments[0] = 0.05
    _set_solid_fill(bg, SOFT)
    bg.line.fill.background()
    bg.shadow.inherit = False
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    x, y, Inches(0.06), h)
    _set_solid_fill(accent, INDIGO)
    tf = bg.text_frame
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, body, size=13, color=SLATE)


def _add_image_card(slide, x, y, w, h, image_path: Path) -> None:
    """White rounded frame around an embedded image."""
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.adjustments[0] = 0.04
    _set_solid_fill(frame, WHITE)
    frame.line.color.rgb = LINE
    frame.line.width = Pt(0.75)
    frame.shadow.inherit = False
    pad = Inches(0.18)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path),
                                 x + pad, y + pad,
                                 w - pad * 2, h - pad * 2)


# ── Slide 1: Title ─────────────────────────────────────────────────────


def slide_title(prs: Presentation, n_pages: int) -> None:
    s = _new_slide(prs)
    # Full-bleed indigo block on the left third with a diagonal accent strip.
    _add_rect(s, 0, 0, Inches(4.4), SLIDE_H, INDIGO_DARK)
    _add_rect(s, Inches(4.4), 0, Inches(0.08), SLIDE_H, ACCENT)
    # Soft circle accent — top-right of left panel.
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.4), Inches(-1.2),
                              Inches(3.2), Inches(3.2))
    _set_solid_fill(circ, INDIGO)
    circ.line.fill.background(); circ.shadow.inherit = False

    # Eyebrow + title on left panel.
    eb = _add_textbox(s, Inches(0.6), Inches(0.7), Inches(3.6), Inches(0.4))
    _set_text(eb, "CIS 2450 · BIG DATA ANALYTICS", size=11,
              bold=True, color=ACCENT)
    title = _add_textbox(s, Inches(0.6), Inches(1.1), Inches(3.6), Inches(2.4))
    _set_text(title, "SEO Ranking\nPredictor", size=44, bold=True, color=WHITE)
    sub = _add_textbox(s, Inches(0.6), Inches(3.0), Inches(3.6), Inches(1.0))
    _set_text(sub, "& Recommendation System", size=20,
              color=RGBColor(0xC7, 0xD2, 0xFE))

    auth = _add_textbox(s, Inches(0.6), Inches(5.3), Inches(3.6), Inches(1.4))
    _add_paragraphs(auth, [
        ("Rahil Patel", {"size": 14, "bold": True, "color": WHITE,
                         "space_after": 2}),
        ("rahilp07@seas.upenn.edu", {"size": 11,
                                     "color": RGBColor(0xC7, 0xD2, 0xFE),
                                     "space_after": 8}),
        ("Ayush Tripathi", {"size": 14, "bold": True, "color": WHITE,
                            "space_after": 2}),
        ("tripath1@seas.upenn.edu", {"size": 11,
                                     "color": RGBColor(0xC7, 0xD2, 0xFE)}),
    ])

    # Right panel: tagline + headline numbers.
    tag = _add_textbox(s, Inches(5.0), Inches(1.4), Inches(7.8), Inches(0.6))
    _set_text(tag, "Predicting Google's top-10 for developer documentation —",
              size=18, color=SLATE_DIM)
    tag2 = _add_textbox(s, Inches(5.0), Inches(1.85), Inches(7.8), Inches(0.6))
    _set_text(tag2, "and explaining every prediction.",
              size=18, bold=True, color=SLATE)

    # Headline stat cards — 4 across.
    df = pd.read_csv(FEATURES_CSV) if FEATURES_CSV.exists() else None
    n_unique = int(df["url"].nunique()) if df is not None else 1297
    n_domains = int(df["domain"].nunique()) if df is not None else 6
    n_features = (df.shape[1] - 5) if df is not None else 72
    cards_y = Inches(3.05)
    card_w = Inches(1.85); card_h = Inches(1.05); card_gap = Inches(0.18)
    base_x = Inches(5.0)
    items = [
        ("Pages", f"{n_unique:,}"),
        ("Domains", f"{n_domains}"),
        ("Features", f"{n_features}"),
        ("Models", "4"),
    ]
    for i, (lab, val) in enumerate(items):
        _add_stat_card(s, base_x + i * (card_w + card_gap), cards_y,
                       card_w, card_h, lab, val)

    pillars = [
        ("Two distinct sources", "scraped HTML + Google SERP rankings"),
        ("Five feature families", "content, metadata, structural, TF-IDF, graph"),
        ("Four-model sweep", "LR → Random Forest → XGBoost → MLP"),
        ("SHAP-driven recommendations", "concrete actions, not vague advice"),
    ]
    for i, (head, body) in enumerate(pillars):
        y = Inches(4.45) + i * Inches(0.55)
        bx = _add_textbox(s, Inches(5.0), y, Inches(8.0), Inches(0.5))
        _add_paragraphs(bx, [
            (head, {"size": 13, "bold": True, "color": INDIGO_DARK}),
        ])
        bx2 = _add_textbox(s, Inches(7.4), y, Inches(5.6), Inches(0.5))
        _add_paragraphs(bx2, [
            (body, {"size": 13, "color": SLATE_DIM}),
        ])


# ── Slide 2: Problem & value proposition ───────────────────────────────


def slide_problem(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="The problem", title="Search ranking is opaque · the inputs are not",
                  page_idx=idx, total_pages=total)

    intro = _add_textbox(s, Inches(0.6), Inches(2.0), Inches(12.2), Inches(0.8))
    _set_text(intro,
              "Documentation authors get no feedback loop on what makes a "
              "page rank. We close that loop by treating SERP placement as a "
              "supervised-learning problem over observable page signals.",
              size=16, color=SLATE_DIM)

    # Three-pillar cards.
    pillars = [
        ("Observable inputs",
         "Content length, heading structure, keyword placement, internal "
         "linking, link-graph authority — all measurable per page."),
        ("Binary classification",
         "`is_top_10` for the topic query derived from each page's "
         "<title>. F1 + PR-AUC (imbalance-aware) as the headline metric."),
        ("Explainable output",
         "SHAP attribution per prediction → a recommendation engine that "
         "tells the author what to change, not just whether they'll rank."),
    ]
    card_w = Inches(4.0); card_h = Inches(3.4); gap = Inches(0.15)
    base_x = Inches(0.6); base_y = Inches(3.2)
    for i, (head, body) in enumerate(pillars):
        x = base_x + i * (card_w + gap)
        # Card body.
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, base_y,
                                  card_w, card_h)
        card.adjustments[0] = 0.04
        _set_solid_fill(card, WHITE)
        card.line.color.rgb = LINE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # Number badge.
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   x + Inches(0.3), base_y + Inches(0.3),
                                   Inches(0.6), Inches(0.6))
        _set_solid_fill(badge, INDIGO)
        badge.line.fill.background(); badge.shadow.inherit = False
        bt = badge.text_frame
        bt.margin_left = bt.margin_right = bt.margin_top = bt.margin_bottom = Inches(0)
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = str(i + 1)
        br.font.name = THEME.font_family; br.font.size = Pt(18)
        br.font.bold = True; br.font.color.rgb = WHITE
        # Heading.
        head_tf = _add_textbox(s, x + Inches(0.3), base_y + Inches(1.05),
                               card_w - Inches(0.6), Inches(0.6))
        _set_text(head_tf, head, size=18, bold=True, color=SLATE)
        # Body.
        body_tf = _add_textbox(s, x + Inches(0.3), base_y + Inches(1.6),
                               card_w - Inches(0.6), card_h - Inches(1.6))
        _set_text(body_tf, body, size=13, color=SLATE_DIM)


# ── Slide 3: Dataset & sources ─────────────────────────────────────────


def slide_dataset(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Dataset", title="Two distinct sources, joined per-page",
                  page_idx=idx, total_pages=total)

    # Left card — Source 1 (HTML).
    sx, sy, sw, sh = Inches(0.6), Inches(2.0), Inches(5.9), Inches(2.6)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, sw, sh)
    card.adjustments[0] = 0.04
    _set_solid_fill(card, WHITE)
    card.line.color.rgb = LINE; card.line.width = Pt(0.75)
    card.shadow.inherit = False
    _add_pill(s, sx + Inches(0.3), sy + Inches(0.3), "Source 1", INDIGO)
    head = _add_textbox(s, sx + Inches(0.3), sy + Inches(0.7), sw - Inches(0.6), Inches(0.5))
    _set_text(head, "Developer-documentation HTML", size=18, bold=True, color=SLATE)
    body = _add_textbox(s, sx + Inches(0.3), sy + Inches(1.2), sw - Inches(0.6), sh - Inches(1.4))
    _add_paragraphs(body, [
        ("Async crawler · robots.txt-aware · 1 req/s per domain",
         {"size": 13, "color": SLATE_DIM, "space_after": 6}),
        ("docs.python.org · developer.mozilla.org · react.dev",
         {"size": 12, "color": SLATE, "space_after": 2}),
        ("nodejs.org · kubernetes.io · fastapi.tiangolo.com",
         {"size": 12, "color": SLATE, "space_after": 8}),
        ("Each page → HTML + JSON sidecar",
         {"size": 11, "color": MUTED, "space_after": 2}),
        ("(URL · fetch_ts · title · outbound_links)",
         {"size": 11, "color": MUTED}),
    ])

    # Right card — Source 2 (SERP).
    sx2 = Inches(6.85)
    card2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx2, sy, sw, sh)
    card2.adjustments[0] = 0.04
    _set_solid_fill(card2, WHITE)
    card2.line.color.rgb = LINE; card2.line.width = Pt(0.75)
    card2.shadow.inherit = False
    _add_pill(s, sx2 + Inches(0.3), sy + Inches(0.3), "Source 2", ACCENT)
    head2 = _add_textbox(s, sx2 + Inches(0.3), sy + Inches(0.7), sw - Inches(0.6), Inches(0.5))
    _set_text(head2, "Google SERP rankings", size=18, bold=True, color=SLATE)
    body2 = _add_textbox(s, sx2 + Inches(0.3), sy + Inches(1.2), sw - Inches(0.6), sh - Inches(1.4))
    _add_paragraphs(body2, [
        ("Brave Search API (free tier) · SerpApi fallback",
         {"size": 13, "color": SLATE_DIM, "space_after": 6}),
        ("One topic query per page → top-10 organic results",
         {"size": 12, "color": SLATE, "space_after": 2}),
        ("Query derived from <title> via boilerplate strip",
         {"size": 12, "color": SLATE, "space_after": 8}),
        ("`is_top_10` = page's own URL appears in its top-10",
         {"size": 11, "color": MUTED}),
    ])

    # Join arrow + result panel.
    join_y = Inches(4.85)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(0.6), join_y, Inches(12.1), Inches(0.55))
    _set_solid_fill(arrow, INDIGO)
    arrow.line.fill.background(); arrow.shadow.inherit = False
    arrow_tf = arrow.text_frame
    arrow_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap = arrow_tf.paragraphs[0]; ap.alignment = PP_ALIGN.CENTER
    arrow_tf.margin_top = arrow_tf.margin_bottom = Inches(0)
    arun = ap.add_run()
    arun.text = "JOIN ON  host + path.lower()"
    arun.font.name = THEME.font_family
    arun.font.size = Pt(13); arun.font.bold = True; arun.font.color.rgb = WHITE

    # Pivot callout — bottom strip.
    _add_callout(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.1),
                 "Per CIS 2450 TA Ricky Gong's 2026-03-29 email, scope was "
                 "narrowed from 50K rows to ~1.3K developer-doc pages — "
                 "free-tier SERP rate-limits made full-scale infeasible "
                 "inside the deadline. Sanctioned, documented in "
                 "data/README.md.")


# ── Slide 4: Feature pipeline ──────────────────────────────────────────


def slide_features(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Feature pipeline",
                  title="72 numeric features across five families",
                  page_idx=idx, total_pages=total)

    families = [
        ("Content", "5", ["text_length", "word_count", "sentence_count",
                          "flesch_reading_ease", "keyword_density"], INDIGO),
        ("Metadata", "4", ["title_length", "has_meta_description",
                           "meta_description_length", "keyword_in_title"], ACCENT),
        ("Structural", "7", ["h1_count", "h2_count", "h3_count",
                             "internal_link_count", "external_link_count",
                             "image_count", "alt_text_coverage"], GOOD),
        ("TF-IDF", "50", ["corpus-fitted top-50 terms",
                          "ngram_range = (1, 2)",
                          "stop_words = english"], INDIGO_DARK),
        ("Graph", "6", ["pagerank (α=0.85)", "hits_hub", "hits_authority",
                        "in_degree", "out_degree", "clustering"], ACCENT),
    ]

    # Five vertical cards across the slide.
    n = len(families)
    total_w = Inches(12.1); gap = Inches(0.18)
    card_w = (total_w - gap * (n - 1)) / n
    base_x = Inches(0.6); base_y = Inches(2.05)
    card_h = Inches(4.2)

    for i, (name, count, items, color) in enumerate(families):
        x = base_x + i * (card_w + gap)
        # Card body.
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, base_y,
                                  card_w, card_h)
        card.adjustments[0] = 0.04
        _set_solid_fill(card, WHITE)
        card.line.color.rgb = LINE; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # Top color bar.
        _add_rect(s, x, base_y, card_w, Inches(0.18), color)
        # Family name.
        head_tf = _add_textbox(s, x + Inches(0.18), base_y + Inches(0.35),
                               card_w - Inches(0.36), Inches(0.4))
        _set_text(head_tf, name, size=15, bold=True, color=SLATE)
        # Count badge — large number.
        cnt_tf = _add_textbox(s, x + Inches(0.18), base_y + Inches(0.75),
                              card_w - Inches(0.36), Inches(0.7))
        _set_text(cnt_tf, count, size=36, bold=True, color=color)
        cnt2 = _add_textbox(s, x + Inches(0.18), base_y + Inches(1.45),
                            card_w - Inches(0.36), Inches(0.3))
        _set_text(cnt2, "columns", size=10, color=MUTED)
        # Item list.
        list_tf = _add_textbox(s, x + Inches(0.18), base_y + Inches(1.85),
                               card_w - Inches(0.36), card_h - Inches(1.95))
        list_tf.word_wrap = True
        for j, it in enumerate(items):
            p = list_tf.paragraphs[0] if j == 0 and not list_tf.paragraphs[0].text \
                else list_tf.add_paragraph()
            p.space_after = Pt(4)
            run = p.add_run(); run.text = it
            run.font.name = "JetBrains Mono"
            run.font.size = Pt(10)
            run.font.color.rgb = SLATE_DIM

    # Bottom callout — joined feature matrix.
    _add_callout(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.85),
                 "All five families are joined per-page on the host+path key "
                 "into data/processed/features.csv — shared across every "
                 "model and the dashboard.")


# ── Slide 5: EDA — class balance + top correlations ────────────────────


def slide_eda_overview(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="EDA · part 1",
                  title="Class balance and feature signal",
                  page_idx=idx, total_pages=total)

    # Two-column layout.
    img_w = Inches(6.0); img_h = Inches(4.2)
    _add_image_card(s, Inches(0.6), Inches(2.0), img_w, img_h,
                    CHART_DIR / "01_class_balance.png")
    _add_image_card(s, Inches(6.85), Inches(2.0), img_w, img_h,
                    CHART_DIR / "03_top_correlations.png")

    # Takeaways under each chart.
    tk1 = _add_textbox(s, Inches(0.6), Inches(6.35), img_w, Inches(0.95))
    _add_paragraphs(tk1, [
        ("Class parity after oversampling", {"size": 13, "bold": True, "color": INDIGO_DARK,
                                              "space_after": 4}),
        ("Minority-class oversampling brought the corpus to a 50/50 split. "
         "PR-AUC stays the headline metric — both classes are equally important.",
         {"size": 12, "color": SLATE_DIM}),
    ])

    tk2 = _add_textbox(s, Inches(6.85), Inches(6.35), img_w, Inches(0.95))
    _add_paragraphs(tk2, [
        ("Strongest signals are structural", {"size": 13, "bold": True, "color": INDIGO_DARK,
                                               "space_after": 4}),
        ("Heading counts, internal links, and graph centrality (PageRank, HITS) "
         "lead the correlation ranking — content length on its own is weaker than expected.",
         {"size": 12, "color": SLATE_DIM}),
    ])


# ── Slide 6: EDA — feature distributions ───────────────────────────────


def slide_eda_distributions(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="EDA · part 2",
                  title="Where the classes diverge",
                  page_idx=idx, total_pages=total)

    img_w = Inches(6.0); img_h = Inches(3.5)
    _add_image_card(s, Inches(0.6), Inches(2.0), img_w, img_h,
                    CHART_DIR / "04_title_length.png")
    _add_image_card(s, Inches(6.85), Inches(2.0), img_w, img_h,
                    CHART_DIR / "05_h2_count.png")

    # Domain breakdown — full width below.
    _add_image_card(s, Inches(0.6), Inches(5.65), Inches(12.25), Inches(1.6),
                    CHART_DIR / "02_domain_breakdown.png")


# ── Slide 7: Link-graph layer ──────────────────────────────────────────


def slide_graph(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Link-graph layer",
                  title="PageRank · HITS · clustering",
                  page_idx=idx, total_pages=total)

    img_w = Inches(6.0); img_h = Inches(3.6)
    _add_image_card(s, Inches(0.6), Inches(2.0), img_w, img_h,
                    CHART_DIR / "06_pagerank_hist.png")
    _add_image_card(s, Inches(6.85), Inches(2.0), img_w, img_h,
                    CHART_DIR / "07_hits_scatter.png")

    # Three method pills.
    methods = [
        ("PageRank", "Stationary distribution · α = 0.85", INDIGO),
        ("HITS", "Hub vs authority · Kleinberg", ACCENT),
        ("Clustering", "Local triangle density per node", GOOD),
    ]
    base_x = Inches(0.6); y = Inches(5.85)
    card_w = Inches(4.0); gap = Inches(0.12); card_h = Inches(1.45)
    for i, (label, body, color) in enumerate(methods):
        x = base_x + i * (card_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.adjustments[0] = 0.05
        _set_solid_fill(card, WHITE)
        card.line.color.rgb = LINE; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        _add_rect(s, x, y, Inches(0.06), card_h, color)
        head_tf = _add_textbox(s, x + Inches(0.25), y + Inches(0.18),
                               card_w - Inches(0.4), Inches(0.45))
        _set_text(head_tf, label, size=15, bold=True, color=color)
        body_tf = _add_textbox(s, x + Inches(0.25), y + Inches(0.65),
                               card_w - Inches(0.4), card_h - Inches(0.7))
        _set_text(body_tf, body, size=12, color=SLATE_DIM)


# ── Slide 8: Modeling sweep ────────────────────────────────────────────


def slide_modeling(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Modeling",
                  title="Baseline → advanced · four-model progression",
                  page_idx=idx, total_pages=total)

    # Four model cards in a row with arrows between them.
    models = [
        ("LR", "Logistic Regression", "Linear baseline · L1 / L2 swept",
         INDIGO),
        ("RF", "Random Forest", "Bagging · n_est, depth, leaf swept",
         ACCENT),
        ("XGB", "XGBoost", "Boosting · η, depth, subsample swept",
         INDIGO_DARK),
        ("MLP", "PyTorch MLP", "2-layer · dropout · batch-norm",
         GOOD),
    ]
    base_x = Inches(0.6); base_y = Inches(2.1)
    card_w = Inches(2.85); card_h = Inches(2.5); arrow_w = Inches(0.2)
    spacing = card_w + arrow_w + Inches(0.05)

    for i, (tag, name, body, color) in enumerate(models):
        x = base_x + i * spacing
        # Card.
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, base_y,
                                  card_w, card_h)
        card.adjustments[0] = 0.05
        _set_solid_fill(card, WHITE)
        card.line.color.rgb = LINE; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # Big tag.
        tag_tf = _add_textbox(s, x + Inches(0.2), base_y + Inches(0.2),
                              card_w - Inches(0.4), Inches(0.7))
        _set_text(tag_tf, tag, size=32, bold=True, color=color)
        # Name.
        name_tf = _add_textbox(s, x + Inches(0.2), base_y + Inches(0.95),
                               card_w - Inches(0.4), Inches(0.5))
        _set_text(name_tf, name, size=15, bold=True, color=SLATE)
        # Description.
        body_tf = _add_textbox(s, x + Inches(0.2), base_y + Inches(1.45),
                               card_w - Inches(0.4), card_h - Inches(1.5))
        _set_text(body_tf, body, size=11, color=SLATE_DIM)
        # Arrow to next card.
        if i < len(models) - 1:
            ax = x + card_w + Inches(0.05)
            ay = base_y + card_h / 2 - Inches(0.15)
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay,
                                    arrow_w, Inches(0.3))
            _set_solid_fill(ar, MUTED)
            ar.line.fill.background(); ar.shadow.inherit = False

    # Below-row callouts: tuner, splitter, imbalance.
    knobs = [
        ("Hyperparameter tuning",
         "RandomizedSearchCV · 30 iter / model · loguniform priors on "
         "continuous knobs · integer ranges on tree depths."),
        ("Cross-validation",
         "Shared StratifiedKFold(n=5, seed=42) across all four models — "
         "comparison-table numbers are apples-to-apples."),
        ("Imbalance handling",
         "class_weight='balanced' (LR/RF) · scale_pos_weight (XGB) · "
         "minority oversampling to class parity for the augmented pass."),
    ]
    cb_y = Inches(4.85); cb_h = Inches(2.4); cb_w = Inches(4.0); cb_gap = Inches(0.18)
    for i, (head, body) in enumerate(knobs):
        x = base_x + i * (cb_w + cb_gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cb_y, cb_w, cb_h)
        card.adjustments[0] = 0.05
        _set_solid_fill(card, SOFT)
        card.line.fill.background(); card.shadow.inherit = False
        head_tf = _add_textbox(s, x + Inches(0.3), cb_y + Inches(0.25),
                               cb_w - Inches(0.6), Inches(0.5))
        _set_text(head_tf, head, size=14, bold=True, color=INDIGO_DARK)
        body_tf = _add_textbox(s, x + Inches(0.3), cb_y + Inches(0.85),
                               cb_w - Inches(0.6), cb_h - Inches(1.0))
        _set_text(body_tf, body, size=11, color=SLATE_DIM)


# ── Slide 9: Results ───────────────────────────────────────────────────


def slide_results(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Results",
                  title="Held-out performance",
                  page_idx=idx, total_pages=total)

    # Big comparison chart on the left.
    _add_image_card(s, Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.5),
                    CHART_DIR / "08_model_comparison.png")

    # Right column: winner stats + confusion thumbnail.
    rx = Inches(8.85); rw = Inches(3.95)

    # Winner banner.
    metrics = _load_metrics()
    winner = _pick_winner(metrics)

    banner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                rx, Inches(2.0), rw, Inches(1.6))
    banner.adjustments[0] = 0.05
    _set_solid_fill(banner, INDIGO_DARK)
    banner.line.fill.background(); banner.shadow.inherit = False

    eb = _add_textbox(s, rx + Inches(0.3), Inches(2.15), rw - Inches(0.6), Inches(0.4))
    _set_text(eb, "WINNER", size=10, bold=True, color=ACCENT)
    name_tf = _add_textbox(s, rx + Inches(0.3), Inches(2.45),
                           rw - Inches(0.6), Inches(0.6))
    _set_text(name_tf, winner["model"], size=24, bold=True, color=WHITE)
    sub_tf = _add_textbox(s, rx + Inches(0.3), Inches(3.0),
                          rw - Inches(0.6), Inches(0.5))
    _set_text(sub_tf,
              f"F1 = {winner['f1']:.3f} · ROC-AUC = {winner['roc_auc']:.3f} · "
              f"PR-AUC = {winner['pr_auc']:.3f}",
              size=12, color=RGBColor(0xC7, 0xD2, 0xFE))

    # Three small metric cards under the banner.
    mcards = [
        ("F1", winner["f1"], INDIGO),
        ("ROC-AUC", winner["roc_auc"], ACCENT),
        ("PR-AUC", winner["pr_auc"], GOOD),
    ]
    mcy = Inches(3.85); mch = Inches(0.95); mcw = (rw - Inches(0.4)) / 3
    for i, (lab, val, color) in enumerate(mcards):
        x = rx + Inches(0.2) + i * (mcw + Inches(0.1))
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, mcy, mcw, mch)
        card.adjustments[0] = 0.08
        _set_solid_fill(card, WHITE)
        card.line.color.rgb = LINE; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        l_tf = _add_textbox(s, x + Inches(0.15), mcy + Inches(0.1),
                            mcw - Inches(0.3), Inches(0.3))
        _set_text(l_tf, lab, size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        v_tf = _add_textbox(s, x + Inches(0.15), mcy + Inches(0.36),
                            mcw - Inches(0.3), Inches(0.5))
        _set_text(v_tf, f"{val:.2f}", size=20, bold=True, color=color,
                  align=PP_ALIGN.CENTER)

    # Confusion matrix preview for the winner.
    cm_path = CHART_DIR / f"09_confusion_{winner['model']}.png"
    if cm_path.exists():
        _add_image_card(s, rx, Inches(4.95), rw, Inches(2.05), cm_path)


def _load_metrics() -> dict[str, dict]:
    out: dict[str, dict] = {}
    if not METRICS_DIR.exists():
        return out
    for path in sorted(METRICS_DIR.glob("*.json")):
        try:
            data = json.loads(path.read_text())
            out[data.get("model", path.stem)] = data
        except Exception:
            continue
    return out


def _pick_winner(metrics: dict[str, dict]) -> dict:
    if not metrics:
        return {"model": "n/a", "f1": 0.0, "roc_auc": 0.0, "pr_auc": 0.0}
    return max(metrics.values(),
               key=lambda m: (m.get("f1", 0.0), m.get("pr_auc", 0.0)))


# ── Slide 10: Dashboard tour ───────────────────────────────────────────


def slide_dashboard(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Dashboard demo",
                  title="Seven tabs · live URL · SHAP-driven recommendations",
                  page_idx=idx, total_pages=total)

    # Mock dashboard "browser chrome" frame on the left (visual placeholder
    # for the actual demo — graders will see the live app during the recording).
    fx, fy, fw, fh = Inches(0.6), Inches(2.0), Inches(7.5), Inches(5.05)
    frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx, fy, fw, fh)
    frame.adjustments[0] = 0.04
    _set_solid_fill(frame, SLATE)
    frame.line.fill.background(); frame.shadow.inherit = False
    # "Sidebar" stripe within the mock.
    _add_rect(s, fx + Inches(0.18), fy + Inches(0.18), Inches(2.0),
              fh - Inches(0.36), INDIGO_DARK)
    # Sidebar title.
    sb_title = _add_textbox(s, fx + Inches(0.3), fy + Inches(0.4),
                            Inches(1.85), Inches(0.4))
    _set_text(sb_title, "SEO RANKING", size=10, bold=True, color=ACCENT)
    sb_sub = _add_textbox(s, fx + Inches(0.3), fy + Inches(0.65),
                          Inches(1.85), Inches(0.35))
    _set_text(sb_sub, "PREDICTOR", size=14, bold=True, color=WHITE)
    # Tab list.
    tabs = ["Predict", "EDA", "Graph", "Models",
            "Recommendations", "What-if", "About"]
    for i, t in enumerate(tabs):
        ty = fy + Inches(1.3) + i * Inches(0.42)
        # Active = first item, with accent.
        is_active = (i == 0)
        if is_active:
            pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      fx + Inches(0.3), ty,
                                      Inches(1.85), Inches(0.32))
            pill.adjustments[0] = 0.4
            _set_solid_fill(pill, ACCENT)
            pill.line.fill.background(); pill.shadow.inherit = False
            tt = pill.text_frame
            tt.margin_top = tt.margin_bottom = Inches(0)
            tt.margin_left = Inches(0.2); tt.margin_right = Inches(0.1)
            tt.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tt.paragraphs[0]
            run = p.add_run(); run.text = t
            run.font.name = THEME.font_family
            run.font.size = Pt(11); run.font.bold = True
            run.font.color.rgb = WHITE
        else:
            tab_tf = _add_textbox(s, fx + Inches(0.45), ty,
                                  Inches(1.7), Inches(0.32))
            _set_text(tab_tf, t, size=11,
                      color=RGBColor(0xC7, 0xD2, 0xFE))
    # Mock content area.
    main_x = fx + Inches(2.35); main_w = fw - Inches(2.55)
    # Banner placeholder (gradient block).
    _add_rect(s, main_x, fy + Inches(0.4), main_w, Inches(0.85),
              INDIGO)
    bn_tf = _add_textbox(s, main_x + Inches(0.25), fy + Inches(0.55),
                         main_w - Inches(0.5), Inches(0.55))
    _set_text(bn_tf, "docs.python.org / library / asyncio.html",
              size=12, bold=True, color=WHITE)
    # Three metric cards mock.
    mc_y = fy + Inches(1.5); mc_h = Inches(1.0)
    mc_w = (main_w - Inches(0.4)) / 3
    mock_metrics = [("XGBoost", "67.2%"), ("Random Forest", "58.4%"),
                    ("Logistic Reg.", "61.1%")]
    for i, (lab, val) in enumerate(mock_metrics):
        x = main_x + Inches(0.05) + i * (mc_w + Inches(0.15))
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, mc_y, mc_w, mc_h)
        c.adjustments[0] = 0.08
        _set_solid_fill(c, WHITE)
        c.line.fill.background(); c.shadow.inherit = False
        ltf = _add_textbox(s, x + Inches(0.12), mc_y + Inches(0.1),
                           mc_w - Inches(0.24), Inches(0.3))
        _set_text(ltf, lab.upper(), size=8, bold=True, color=MUTED)
        vtf = _add_textbox(s, x + Inches(0.12), mc_y + Inches(0.32),
                           mc_w - Inches(0.24), Inches(0.55))
        _set_text(vtf, val, size=20, bold=True, color=GOOD)
    # SHAP bars mock.
    shap_y = fy + Inches(2.85)
    shap_tf = _add_textbox(s, main_x + Inches(0.05), shap_y,
                           main_w - Inches(0.1), Inches(0.35))
    _set_text(shap_tf, "Per-prediction SHAP attribution", size=11,
              bold=True, color=WHITE)
    bar_specs = [("h2_count", 0.78, GOOD), ("title_length", 0.45, GOOD),
                 ("keyword_density", -0.32, BG), ("pagerank", 0.55, GOOD),
                 ("internal_link_count", -0.18, BG)]
    bar_y0 = shap_y + Inches(0.4)
    name_w = Inches(1.6); bar_total_w = main_w - name_w - Inches(0.2)
    for i, (name, val, color) in enumerate(bar_specs):
        y = bar_y0 + i * Inches(0.32)
        nm_tf = _add_textbox(s, main_x + Inches(0.05), y, name_w, Inches(0.3))
        _set_text(nm_tf, name, size=9, color=RGBColor(0xC7, 0xD2, 0xFE),
                  family="JetBrains Mono")
        # Track.
        track = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   main_x + name_w + Inches(0.05),
                                   y + Inches(0.07),
                                   bar_total_w, Inches(0.14))
        track.adjustments[0] = 0.5
        _set_solid_fill(track, RGBColor(0x1E, 0x29, 0x3B))
        track.line.fill.background(); track.shadow.inherit = False
        # Filled portion.
        mag = abs(val)
        center_x = main_x + name_w + Inches(0.05) + bar_total_w / 2
        fill_w = bar_total_w * mag / 2
        fill_color = GOOD if val >= 0 else BAD
        if val >= 0:
            fill_x = center_x
        else:
            fill_x = center_x - fill_w
        fill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  fill_x, y + Inches(0.07),
                                  fill_w, Inches(0.14))
        fill.adjustments[0] = 0.5
        _set_solid_fill(fill, fill_color)
        fill.line.fill.background(); fill.shadow.inherit = False

    # Right column — feature list of the dashboard.
    rx = Inches(8.4); rw = Inches(4.45)
    sec_tf = _add_textbox(s, rx, Inches(2.0), rw, Inches(0.4))
    _set_text(sec_tf, "What the dashboard does", size=14, bold=True,
              color=INDIGO_DARK)
    feats = [
        ("Predict", "Live URL → top-10 probability per loaded model"),
        ("EDA", "Class balance · correlations · distribution explorer"),
        ("Graph", "PageRank · HITS scatter · URL hierarchy network"),
        ("Models", "Metrics table · ROC + PR curves · confusion matrices"),
        ("Recommendations", "Hybrid rule + SHAP-ranked actions per page"),
        ("What-if", "Slider counterfactuals · live probability deltas"),
    ]
    fy_base = Inches(2.55)
    for i, (head, body) in enumerate(feats):
        y = fy_base + i * Inches(0.7)
        # Tab pill.
        _add_pill(s, rx, y, head, INDIGO if i % 2 == 0 else ACCENT,
                  size=10)
        b_tf = _add_textbox(s, rx, y + Inches(0.36), rw, Inches(0.35))
        _set_text(b_tf, body, size=11, color=SLATE_DIM)


# ── Slide 11: Insights & implications ──────────────────────────────────


def slide_insights(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Insights",
                  title="What the model actually learned",
                  page_idx=idx, total_pages=total)

    # Three-column "lesson" cards.
    items = [
        ("Structure beats length",
         "Heading-count and link-graph features outrank raw word count "
         "in the importance ranking. Authors should optimise H2 / H3 "
         "structure before pushing for longer pages."),
        ("Authority is observable",
         "PageRank and HITS authority correlate with `is_top_10` more "
         "strongly than any single content feature. Internal-link "
         "topology — what we already control on our own site — is "
         "actionable signal."),
        ("Title craft is high-leverage",
         "Both `keyword_in_title` and `title_length` show clean class "
         "separation. The What-if simulator confirms: changing title "
         "length by ±15 chars moves predicted probability by 6-12 pts."),
    ]
    base_x = Inches(0.6); base_y = Inches(2.1)
    card_w = Inches(4.0); card_h = Inches(3.6); gap = Inches(0.18)
    for i, (head, body) in enumerate(items):
        x = base_x + i * (card_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, base_y,
                                  card_w, card_h)
        card.adjustments[0] = 0.04
        _set_solid_fill(card, WHITE)
        card.line.color.rgb = LINE; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        # Big numeric badge.
        n_tf = _add_textbox(s, x + Inches(0.3), base_y + Inches(0.2),
                            card_w - Inches(0.6), Inches(0.7))
        _set_text(n_tf, f"0{i + 1}", size=32, bold=True, color=INDIGO)
        head_tf = _add_textbox(s, x + Inches(0.3), base_y + Inches(1.05),
                               card_w - Inches(0.6), Inches(0.6))
        _set_text(head_tf, head, size=17, bold=True, color=SLATE)
        body_tf = _add_textbox(s, x + Inches(0.3), base_y + Inches(1.7),
                               card_w - Inches(0.6), card_h - Inches(1.85))
        _set_text(body_tf, body, size=13, color=SLATE_DIM)

    _add_callout(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.3),
                 "Practical implication: the recommendation engine prioritises "
                 "structural fixes (heading counts, alt-text coverage, "
                 "internal-link distribution) ahead of content rewrites — "
                 "lower edit cost, higher predicted lift.")


# ── Slide 12: Challenges, limitations, future work ─────────────────────


def slide_challenges(prs: Presentation, idx: int, total: int) -> None:
    s = _new_slide(prs)
    _slide_chrome(s, eyebrow="Challenges & future work",
                  title="What we'd ship next",
                  page_idx=idx, total_pages=total)

    # Two-column layout: challenges (left) vs future (right).
    base_y = Inches(2.0)
    col_w = Inches(6.0); col_gap = Inches(0.25); col_h = Inches(5.1)

    # Left: challenges.
    lx = Inches(0.6)
    lcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lx, base_y,
                               col_w, col_h)
    lcard.adjustments[0] = 0.04
    _set_solid_fill(lcard, WHITE)
    lcard.line.color.rgb = LINE; lcard.line.width = Pt(0.75)
    lcard.shadow.inherit = False
    _add_pill(s, lx + Inches(0.3), base_y + Inches(0.3), "Challenges", BAD)
    challenges = [
        ("Free-tier API limits",
         "Brave Search caps at ~2K queries/month — bounded the corpus "
         "to ~1.3K pages."),
        ("Class imbalance in raw labels",
         "Mitigated with class_weight + minority oversampling to parity."),
        ("Live URLs lack graph signal",
         "Page isn't in the training graph; we median-fill its graph "
         "features rather than zero-fill (zero biases predictions down)."),
        ("Title-derived queries can mislead",
         "Some <title> tags are version numbers or 'Overview'. The "
         "dashboard exposes a manual query override for evaluation."),
    ]
    for i, (head, body) in enumerate(challenges):
        y = base_y + Inches(0.85) + i * Inches(1.05)
        head_tf = _add_textbox(s, lx + Inches(0.3), y, col_w - Inches(0.6),
                               Inches(0.4))
        _set_text(head_tf, head, size=14, bold=True, color=SLATE)
        body_tf = _add_textbox(s, lx + Inches(0.3), y + Inches(0.4),
                               col_w - Inches(0.6), Inches(0.65))
        _set_text(body_tf, body, size=12, color=SLATE_DIM)

    # Right: future work.
    rx = lx + col_w + col_gap
    rcard = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, base_y,
                               col_w, col_h)
    rcard.adjustments[0] = 0.04
    _set_solid_fill(rcard, WHITE)
    rcard.line.color.rgb = LINE; rcard.line.width = Pt(0.75)
    rcard.shadow.inherit = False
    _add_pill(s, rx + Inches(0.3), base_y + Inches(0.3), "Future work", GOOD)
    futures = [
        ("Per-domain calibrated thresholds",
         "0.5 cutoff is naive — real SERPs differ in expected positive "
         "rate per domain. Calibrate via Platt / isotonic on a held-out "
         "validation slice per host."),
        ("Periodic re-scrape + drift monitor",
         "SERP rankings shift weekly. Schedule a refresh, log feature "
         "drift between snapshots, retrigger training when drift > τ."),
        ("Recommendation A/B harness",
         "Apply a recommended change to a copy of a page, re-scrape, "
         "compare actual rank movement against predicted lift — closes "
         "the loop on whether the model is causal or merely correlated."),
        ("Cross-engine generalisation",
         "Today: Google via Brave / SerpApi. Add Bing / DuckDuckGo "
         "labels — does the same model rank well across engines?"),
    ]
    for i, (head, body) in enumerate(futures):
        y = base_y + Inches(0.85) + i * Inches(1.05)
        head_tf = _add_textbox(s, rx + Inches(0.3), y, col_w - Inches(0.6),
                               Inches(0.4))
        _set_text(head_tf, head, size=14, bold=True, color=SLATE)
        body_tf = _add_textbox(s, rx + Inches(0.3), y + Inches(0.4),
                               col_w - Inches(0.6), Inches(0.65))
        _set_text(body_tf, body, size=12, color=SLATE_DIM)


# ── Slide 13: Thank you / Q&A ──────────────────────────────────────────


def slide_thanks(prs: Presentation) -> None:
    s = _new_slide(prs)
    # Full-bleed indigo panel.
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INDIGO_DARK)
    # Diagonal accent strip.
    _add_rect(s, 0, Inches(7.0), SLIDE_W, Inches(0.5), ACCENT)
    # Big circle accent.
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2.5),
                              Inches(7), Inches(7))
    _set_solid_fill(circ, INDIGO)
    circ.line.fill.background(); circ.shadow.inherit = False

    eb = _add_textbox(s, Inches(0.8), Inches(2.4), Inches(8), Inches(0.6))
    _set_text(eb, "QUESTIONS?", size=14, bold=True, color=ACCENT)

    title = _add_textbox(s, Inches(0.8), Inches(2.9), Inches(11), Inches(1.6))
    _set_text(title, "Thank you.", size=84, bold=True, color=WHITE)

    sub = _add_textbox(s, Inches(0.8), Inches(4.7), Inches(11), Inches(0.7))
    _set_text(sub, "SEO Ranking Predictor · CIS 2450 Final Project",
              size=20, color=RGBColor(0xC7, 0xD2, 0xFE))

    auth = _add_textbox(s, Inches(0.8), Inches(5.4), Inches(11), Inches(0.5))
    _set_text(auth, "Rahil Patel  ·  Ayush Tripathi",
              size=16, color=RGBColor(0xC7, 0xD2, 0xFE))


# ── Main ───────────────────────────────────────────────────────────────


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_dataset,
        slide_features,
        slide_eda_overview,
        slide_eda_distributions,
        slide_graph,
        slide_modeling,
        slide_results,
        slide_dashboard,
        slide_insights,
        slide_challenges,
    ]
    total = len(builders) + 1  # +1 for thanks slide

    # First (title) takes only `prs, total` (different signature); rest take
    # `prs, idx, total`.
    slide_title(prs, total)
    for i, fn in enumerate(builders[1:], start=2):
        fn(prs, i, total)
    slide_thanks(prs)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"saved → {OUT_PPTX} ({OUT_PPTX.stat().st_size // 1024} KB, "
          f"{total} slides)")


if __name__ == "__main__":
    main()






